"""Update the PowerPoint deck with corrected text and regenerated visuals."""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DECK_PATH = ROOT / "deliverables" / "montgomery_crash_safety_presentation.pptx"


TEXT_REPLACEMENTS = {
    "As Resources Are Limited So As People’s Attention": "Limited resources require focused safety investment",
    "Now We Need to Consider (Equity & Vulnerable Users)": "Priority scoring adds equity and vulnerable users",
    "Top 10 Segments by composite score": "Top 10 Segments by CEI-aware priority score",
}


SLIDE_IMAGES = {
    6: "assets/chart_01_annual_crashes.png",
    9: "assets/chart_03_vru_dui_ksi_share.png",
    11: "assets/chart_04_average_cost.png",
    12: "assets/chart_03_vru_dui_ksi_share.png",
    15: "assets/chart_05_hin_map.png",
    16: "assets/chart_06_hin_concentration.png",
    17: "assets/chart_07_top_ksi_segments.png",
    19: "assets/chart_08_exposure_rate.png",
    23: "assets/chart_09_priority_score.png",
    24: "assets/chart_10_cei_priority_map.png",
}


def replace_text(prs: Presentation) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    for old, new in TEXT_REPLACEMENTS.items():
                        text = text.replace(old, new)
                    run.text = text


def delete_pictures(slide) -> None:
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape._element.getparent().remove(shape._element)


def add_picture_fit(slide, image_path: Path, left: int, top: int, width: int, height: int) -> None:
    with Image.open(image_path) as image:
        image_width, image_height = image.size

    image_ratio = image_width / image_height
    box_ratio = width / height
    if image_ratio >= box_ratio:
        fitted_width = width
        fitted_height = round(width / image_ratio)
    else:
        fitted_height = height
        fitted_width = round(height * image_ratio)

    fitted_left = left + round((width - fitted_width) / 2)
    fitted_top = top + round((height - fitted_height) / 2)
    slide.shapes.add_picture(
        str(image_path),
        fitted_left,
        fitted_top,
        width=fitted_width,
        height=fitted_height,
    )


def add_chart(slide, image_path: Path, prs: Presentation, full: bool = False) -> None:
    if full:
        add_picture_fit(
            slide,
            image_path,
            left=Inches(0.35),
            top=Inches(0.8),
            width=prs.slide_width - Inches(0.7),
            height=prs.slide_height - Inches(1.05),
        )
        return
    add_picture_fit(
        slide,
        image_path,
        left=Inches(0.55),
        top=Inches(1.05),
        width=prs.slide_width - Inches(1.1),
        height=prs.slide_height - Inches(1.35),
    )


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def add_or_update_title(slide, text: str, prs: Presentation) -> None:
    title = slide.shapes.title
    if title is not None:
        title.text = text
        title.left = Inches(0.55)
        title.top = Inches(0.25)
        title.width = prs.slide_width - Inches(1.1)
        title.height = Inches(0.55)
        return

    candidates = [
        shape
        for shape in list(slide.shapes)
        if getattr(shape, "has_text_frame", False)
        and ("CEI disadvantaged tracts" in shape.text or "priority corridors" in shape.text.lower())
    ]
    if candidates:
        box = candidates[0]
        for duplicate in candidates[1:]:
            remove_shape(duplicate)
        box.left = Inches(0.55)
        box.top = Inches(0.25)
        box.width = prs.slide_width - Inches(1.1)
        box.height = Inches(0.55)
    else:
        box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), prs.slide_width - Inches(1.1), Inches(0.55))

    box.text_frame.text = text
    for paragraph in box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(24)
            run.font.bold = True


def clamp_shapes_to_slide(prs: Presentation) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.width > prs.slide_width:
                shape.width = prs.slide_width
            if shape.height > prs.slide_height:
                shape.height = prs.slide_height
            if shape.left < 0:
                shape.left = 0
            if shape.top < 0:
                shape.top = 0
            if shape.left + shape.width > prs.slide_width:
                shape.left = max(0, prs.slide_width - shape.width)
            if shape.top + shape.height > prs.slide_height:
                shape.top = max(0, prs.slide_height - shape.height)


def update_images(prs: Presentation) -> None:
    for slide_number, rel_path in SLIDE_IMAGES.items():
        slide = prs.slides[slide_number - 1]
        delete_pictures(slide)
        image_path = ROOT / rel_path
        add_chart(slide, image_path, prs, full=slide_number in {15, 24})

    add_or_update_title(prs.slides[23], "CEI disadvantaged tracts and priority corridors", prs)


def update_key_slide_text(prs: Presentation) -> None:
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if getattr(shape, "has_text_frame", False) and "From" in shape.text:
            shape.text = (
                "From 2020-2024, Montgomery County recorded 59,476 crashes. "
                "Of these, 1,400 were fatal or suspected serious injury crashes (KSI), "
                "with an estimated $7.1 billion in societal cost. To reach Vision Zero, "
                "the County should focus engineering, enforcement, and policy on the "
                "corridors and conditions most strongly linked to severe crashes."
            )

    slide22 = prs.slides[21]
    replacements_by_text = {
        "KSI crash burden": "KSI crash burden",
        "(0.5 x Fatal collisions)": "(1.00 x KSI burden)",
        "(0.25 x Vulnerable Users collisions)": "(0.50 x VRU KSI)",
        "(0.25 x KSI crash burden\nIn disadvantaged areas)": "(0.25 x KSI in disadvantaged tracts)",
        "(0.25 x exposure-adjusted risk)": "(0.25 x exposure-adjusted risk)",
    }
    for shape in slide22.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text in replacements_by_text:
                shape.text = replacements_by_text[text]

    slide26 = prs.slides[25]
    for shape in slide26.shapes:
        if getattr(shape, "has_text_frame", False) and "Re-score" in shape.text:
            shape.text = (
                "Prioritizes high-harm corridors, including corridors in disadvantaged tracts\n\n"
                "Re-score every 2-3 years\n\n"
                "Phased delivery:\n"
                "Year 1: quick wins such as LPIs, daylighting, and pilot road diets\n"
                "Years 1-3: corridor redesigns\n"
                "Years 3-5: major reconstructions\n\n"
                "Track progress: KSI, HIN share, VRU KSI, CEI contribution, and exposure-adjusted risk"
            )


def main() -> None:
    prs = Presentation(DECK_PATH)
    replace_text(prs)
    update_key_slide_text(prs)
    update_images(prs)
    clamp_shapes_to_slide(prs)
    prs.save(DECK_PATH)
    print(f"Updated {DECK_PATH}")


if __name__ == "__main__":
    main()
